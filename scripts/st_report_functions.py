import numpy as np
from math import ceil
import streamlit as st
import pandas as pd
import math
import plotly.graph_objects as go
# reset Plotly theme after streamlit import
import plotly.io as pio
pio.templates.default = 'plotly'
from colour import Color

from util import filter_dict_by_string, run_script

from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor

from io import BytesIO

from rdkit import Chem
from rdkit.Chem import AllChem

from PIL import Image



legend_font_size = 22
title_font_size = 22

                     # River Blue           Tuck Orange         Rich Spring Green    Web Violet            Bonfire Orange
dm_tertiary_colors = ['rgb(38, 122, 186)', 'rgb(227, 45, 28)', 'rgb(165, 215, 95)', 'rgb(138, 105, 150)', 'rgb(255, 160, 15)']
tab20_colormap = ['rgb(31, 119, 180)', 'rgb(174, 199, 232)', 'rgb(255, 127, 14)', 'rgb(255, 187, 120)',
                  'rgb(44, 160, 44)', 'rgb(152, 223, 138)', 'rgb(214, 39, 40)', 'rgb(255, 152, 150)',
                  'rgb(148, 103, 189)', 'rgb(197, 176, 213)', 'rgb(140, 86, 75)', 'rgb(196, 156, 148)',
                  'rgb(227, 119, 194)', 'rgb(247, 182, 210)', 'rgb(127, 127, 127)', 'rgb(199, 199, 199)',
                  'rgb(188, 189, 34)', 'rgb(219, 219, 141)', 'rgb(23, 190, 207)', 'rgb(158, 218, 229)']

def make_dssp_plot(dataframe_subset, data_type, report=False):
    fig = go.Figure()
    color_counter = 0
    if report is True:
        color_palette = tab20_colormap
    else:
        color_palette = dm_tertiary_colors
    for name, dataset in dataframe_subset.items():
        info = name.split('|')
        for data in data_type:
            fig.add_trace(go.Scatter(
                y=dataset[data],
                x=dataset.index.str.replace('_', ' '),
                # name=f'Replica{info[1]}'
                name=f'{data.replace("DSSP_", "")}',
                line=dict(color=color_palette[color_counter])
            ))
            color_counter += 1
    
    fig.update_layout(
        # title = f'{info[0]} {data_type.replace("_", " ")}'
        title_text=f'Secondary Structure',
        font=dict(size=title_font_size),
        plot_bgcolor='rgb(255,255,255)',
        legend=dict(
            orientation = 'h',
            yanchor='bottom',
            y=0.9,
            xanchor='right',
            x=1,
            font=dict(size=legend_font_size)
            )
    )

    fig.update_yaxes(
        title='Secondary Structure Fraction',
        range=[0,1],
        showgrid=True,
        gridcolor='rgb(226, 226, 226)',
        showline=True,
        linecolor='black',
        linewidth=2,
        tickfont=dict(size=18),
        title_font=dict(size=22),
    )
    fig.update_xaxes(
        title='',
        showline=True,
        linecolor='black',
        linewidth=2,
        tickfont=dict(size=18),
        title_font=dict(size=22), 
    )

    if report is True:
        fig.update_layout(title_text=f'{info[0].replace("_", " ")} {info[2].replace("_", " ")}')
        for trace in fig.data:
            trace.showlegend = False

    return fig


def make_salpha_plot(dataframe_subset):
    fig = go.Figure()
    for name, dataset in dataframe_subset.items():
        info = name.split('|')
        fig.add_trace(go.Scatter(
            y=dataset['salpha_free_energy'],
            x=dataset['salpha_edges'],
            name=f'Replica{info[1]}'
        ))
    fig.update_layout(
        # title = f'{info[0]} {data_type.replace("_", " ")}'
        title_text=f'',
        font=dict(size=title_font_size),
        plot_bgcolor='rgb(255,255,255)',
        legend=dict(
            orientation = 'h',
            yanchor='bottom',
            y=0.9,
            xanchor='right',
            x=1,
            font=dict(size=legend_font_size)
            )
    )

    fig.update_yaxes(
        title='Free Energy (KCal/mol)',
        showgrid=True,
        gridcolor='rgb(226, 226, 226)',
        showline=True,
        linecolor='black',
        linewidth=2,
        tickfont=dict(size=18),
        title_font=dict(size=22),
    )
    fig.update_xaxes(
        title='Sα',
        showline=True,
        linecolor='black',
        linewidth=2,
        tickfont=dict(size=18),
        title_font=dict(size=22), 
    )

    return fig


def make_contact_probability_plot(dataframe_subset, contact_probability_type, title_plot='Ligand Contact Probability', show_legend=False):
    contact_probability = go.Figure()
    color_counter = 0

    # for name, dataset in dataframe_subset.items():
    # info = name.split('|')
    for data_type in contact_probability_type:
        try:
            contact_probability.add_trace(go.Scatter(
                y=dataframe_subset[data_type] + dataframe_subset[f'{data_type}_error'],
                x=dataframe_subset.index.str.replace('_', ' '),
                showlegend=False,
                mode='lines',
                line=dict(width=0),
                name='',
            ))
            contact_probability.add_trace(go.Scatter(
                y=dataframe_subset[data_type] - dataframe_subset[f'{data_type}_error'],
                x=dataframe_subset.index.str.replace('_', ' '),
                showlegend=False,
                mode='lines',
                line=dict(width=0),
                fillcolor='rgba(226, 226, 226, 0.3)',
                fill='tonexty',
                name='',
            ))
        except: pass
        contact_probability.add_trace(go.Scatter(
            y=dataframe_subset[data_type],
            x=dataframe_subset.index.str.replace('_', ' '),
            name=data_type.replace("_", " "),
            showlegend=show_legend,
            line=dict(color=dm_tertiary_colors[color_counter])
        ))
        color_counter += 1
    contact_probability.update_yaxes(
        title='Contact probability',
        # range=[0,1],
        showgrid=True,
        gridwidth=1,
        gridcolor='rgb(226, 226, 226)',
        showline=True,
        linecolor='black',
        linewidth=2,
        tickfont=dict(size=18),
        title_font=dict(size=22),
    )
    contact_probability.update_xaxes(
        title='',
        showline=True,
        linecolor='black',
        linewidth=2,
        tickfont=dict(size=18),
        title_font=dict(size=22),
    )
    contact_probability.update_layout(
        title_text=title_plot,
        font=dict(size=title_font_size),
        plot_bgcolor='rgb(255,255,255)',
        legend=dict(
            orientation = 'h',
            yanchor='bottom',
            y=0.9,
            xanchor='right',
            x=1,
            font=dict(size=legend_font_size)
        )
    )

    return contact_probability


def make_kd_plot(dataframe_subset):
    kd_plot = go.Figure()
    # for name, dataset in dataframe_subset.items():
    # info = name.split('|')
    kd_plot.add_trace(go.Scatter(
        y=dataframe_subset['Kd_error_up'],
        x=dataframe_subset.index/1000/1000,
        showlegend=False,
        mode='lines',
        line=dict(width=0),
        name='',
        # name=f'Replica{info[1]}'
    ))
    kd_plot.add_trace(go.Scatter(
        y=dataframe_subset['Kd_error_low'],
        x=dataframe_subset.index/1000/1000,
        showlegend=False,
        mode='lines',
        line=dict(width=0),
        fillcolor='rgba(226, 226, 226, 0.3)',
        fill='tonexty',
        name='',
    ))
    kd_plot.add_trace(go.Scatter(
        y=dataframe_subset['Kd'],
        x=dataframe_subset.index/1000/1000,
        showlegend=False,
        # name='Kd',
        line=dict(color=dm_tertiary_colors[0])
    ))
    values_without_inf = [x for x in dataframe_subset['Kd'].to_list() if not math.isinf(x)]
    kd_plot.update_yaxes(
        title='Bound Fraction',
        range=[0, 1],
        # range=[0, max(values_without_inf)+2],
        # range=[0, dataframe_subset['Kd'].max()+2],
        showgrid=True,
        gridwidth=1,
        gridcolor='rgb(226, 226, 226)',
        showline=True,
        linecolor='black',
        linewidth=2,
        tickfont=dict(size=18),
        title_font=dict(size=22),
    )
    kd_plot.update_xaxes(
        title='Time (us)',
        showline=True,
        linecolor='black',
        linewidth=2,
        tickfont=dict(size=18),
        title_font=dict(size=22),
    )
    kd_plot.update_layout(
        title_text=f'',
        # title_text=f'Simulated Kd',
        font=dict(size=title_font_size),
        plot_bgcolor='rgb(255,255,255)',
    )

    return kd_plot


def make_contact_matrix_plot(dataframe_subset, dual='', report=False):
    contact_matrix_fig = go.Figure()
    # This for loop can actually be removed.
    for name, dataset in dataframe_subset.items():
        info = name.split('|')
        contact_matrix_fig.add_trace(go.Heatmap(
            z=dataset,
            x=dataset.index.str.replace('_', ' '),
            y=dataset.index.str.replace('_', ' '),
            colorscale='Jet'
        ))
    contact_matrix_fig.update_yaxes(
        tickfont=dict(size=18),
        title_font=dict(size=22),
    )
    contact_matrix_fig.update_xaxes(
        tickfont=dict(size=18),
        title_font=dict(size=22),
    )
    if report is True:
        contact_matrix_fig.update_layout(
            title_text=f'{info[0]} {info[1]} {info[2].replace("_", " ")}',
            font=dict(size=title_font_size),
        )
    else:
        contact_matrix_fig.update_layout(
            title_text=f'{dual}Contact Probability',
            font=dict(size=title_font_size),
        )

    return contact_matrix_fig


def make_fes_gyr_salpha_plot(dataframe_subset, report=False):
    fig = go.Figure()
    for k, v in dataframe_subset.items():
        info = k.split('|')
        fig.add_trace(go.Contour(
            z=np.flip(
                v.to_numpy(),
                axis=0
                ),
            x=list(v.index),
            y=list(v.columns),
            zmin=0.1,
            zmax=3,
            line_smoothing = 0.5,
            ncontours=20,
            contours_coloring='heatmap',
            colorscale='Jet')
        )
    fig.update_yaxes(
        title='Radius of Gyration (nm)',
        tickfont=dict(size=18),
        title_font=dict(size=22),
    )
    fig.update_xaxes(
        title='Sα',
        tickfont=dict(size=18),
        title_font=dict(size=22),
    )
    if report is True:
        fig.update_layout(
            title_text=f'{info[0]} {info[1]} {info[2].replace("_", " ")}',
            font=dict(size=title_font_size),
        )
    else:
        fig.update_layout(
            title_text=f'Free Energy (KCal/mol)',
            font=dict(size=title_font_size),
        )
    return fig


def make_ligand_contacts_comparison(hphob_contact_probability_fig, aromatic_contact_probability_fig, hbond_contact_probability_fig=None):

    contact_probabilities_comparison_fig = go.Figure()
    selected_trace = list(hphob_contact_probability_fig.select_traces(selector=dict(name='hydrophobic contacts')))
    contact_probabilities_comparison_fig.add_trace(selected_trace[0])
    contact_probabilities_comparison_fig.update_traces(selector=dict(name='hydrophobic contacts'), line=dict(color=dm_tertiary_colors[1]))


    selected_trace = list(aromatic_contact_probability_fig.select_traces(selector=dict(name='aromatic stacking')))
    contact_probabilities_comparison_fig.add_trace(selected_trace[0])
    contact_probabilities_comparison_fig.update_traces(selector=dict(name='aromatic stacking'), line=dict(color=dm_tertiary_colors[0]))
 
 
    if hbond_contact_probability_fig is not None:
        selected_trace = list(hbond_contact_probability_fig.select_traces(selector=dict(name='Hbonds average')))
        contact_probabilities_comparison_fig.add_trace(selected_trace[0])    
        contact_probabilities_comparison_fig.update_traces(selector=dict(name='Hbonds average'), line=dict(color=dm_tertiary_colors[2]))

    for trace in contact_probabilities_comparison_fig.data:
        trace.showlegend = True
    
    contact_probabilities_comparison_fig.update_yaxes(
        title='Contact probability',
        # range=[0,1],
        showgrid=True,
        gridwidth=1,
        gridcolor='rgb(226, 226, 226)',
        showline=True,
        linecolor='black',
        linewidth=2,
        tickfont=dict(size=18),
        title_font=dict(size=22),
    )
    contact_probabilities_comparison_fig.update_xaxes(
        title='',
        showline=True,
        linecolor='black',
        linewidth=2,
        tickfont=dict(size=18),
        title_font=dict(size=22),
    )
    contact_probabilities_comparison_fig.update_layout(
        # title_text=f'Ligand Contact Probability',
        font=dict(size=title_font_size),
        plot_bgcolor='rgb(255,255,255)',
        legend=dict(
            orientation = 'h',
            yanchor='bottom',
            y=0.9,
            xanchor='right',
            x=1,
            font=dict(size=legend_font_size)
        )
    )
    return contact_probabilities_comparison_fig


def make_pymol_images(title, contact_dataframe, residue_name, simulation_output_folder, max_probability=None, min_probability=None, min_color='pink', max_color='blue', image_output_folder=None):
    
    for simulation, interaction_dataframe in contact_dataframe.items():
        name, replica, trajectory, data = simulation.split('|')
        
        
        # Make 2D coordinates for the image using rdkit

        rdkit_molecule = Chem.MolFromPDBFile(f'{simulation_output_folder}/{replica}/ligand_{replica}.pdb', removeHs=False)
        AllChem.Compute2DCoords(rdkit_molecule)
        Chem.MolToPDBFile(rdkit_molecule, f'{simulation_output_folder}/{replica}/ligand_{replica}_2D.pdb')

        
        interaction_dataframe['Probability'] = interaction_dataframe['Probability'].round(3)
        if max_probability is None:
            max_probability = interaction_dataframe['Probability'].max()
        if min_probability is None:
            min_probability = interaction_dataframe['Probability'][interaction_dataframe['Probability'] > 0 ].min()
        probability_range = np.arange(min_probability, max_probability+0.0004, step=0.0001) # +0.0004 for some reason I needed to add this
        color_spectrum = list(Color(min_color).range_to(Color(max_color), len(probability_range)))
        colorbar_color_spectrum = list(Color(min_color).range_to(Color(max_color), 1000))
    
        new_colorbar_color_spectrum = []
        for color in colorbar_color_spectrum:
            rgb_color = list(color.rgb)
            new_rgb_color = [int(c*255) for c in rgb_color]
            new_colorbar_color_spectrum.append(new_rgb_color)
    
        color_dictionary_to_map = {}
        for value, color in zip(probability_range, color_spectrum):
            value = np.round(value, 4)
            color_dictionary_to_map[value] = color.hex_l
            # color_dictionary_to_map[value] = color.rgb

        interaction_dataframe['colors'] = interaction_dataframe['Probability'].map(color_dictionary_to_map)
        
        # x_values = [ 1 for e in new_colorbar_color_spectrum] 
        x_labels = ['' for e in new_colorbar_color_spectrum]
        x_labels[0] = min_probability
        x_labels[-1] = max_probability
    
        color_bar = go.Figure()
        color_bar.add_trace(go.Image(z=[new_colorbar_color_spectrum], dy=30))
        color_bar.update_yaxes(showticklabels=False)
        # color_bar.update_xaxes(tickangle=90, tickvals = list(range(0, len(x_labels))), ticktext=x_labels, tickfont=dict(size=38))
        color_bar.update_xaxes(showticklabels=False)
        color_bar.update_layout(xaxis={'automargin':True})
        color_bar.update_layout(
            margin=dict(r=150),
            # ^^ making a bit of space for the annotation
            annotations=[
                dict(
                    text=f"{title} {interaction_dataframe['Probability'].max()}-{interaction_dataframe['Probability'][interaction_dataframe['Probability'] > 0 ].min()}",
                    font_size=45,
                    textangle=180,
                    showarrow=False,
                    # ^^ appearance
                    xref="paper",
                    yref="paper",
                    x=0.5,
                    y=0.25,
                    # ^^ position
                )
            ]
        )

        color_bar.write_image(f'{image_output_folder}/{name}_{replica}_{trajectory}_{data}_colorbar.png', format='png', scale=6)

        with open(f'{simulation_output_folder}/{replica}/{name}_{replica}_{trajectory}_{data}.pml', 'w') as file:
            file.write(f"cmd.load('{simulation_output_folder}/{replica}/ligand_{replica}_2D.pdb', 'structure')\n")
            # file.write(f"cmd.delete('structure not {selection}')\n")
            # file.write(f"cmd.remove('all not {selection}')\n")
            # file.write(f"cmd.load('{trajectory}', 'structure')\n")
            file.write(f"cmd.center('{residue_name}')\n")
            file.write(f"cmd.intra_fit('{residue_name}', 0)\n")
            file.write(f"cmd.hide('', '')\n")
            # file.write(f"cmd.hide('', 'name NA')\n")
            # file.write(f"cmd.hide('', 'name CL')\n")
            file.write(f"cmd.show('sticks', '{residue_name}')\n")
            # file.write(f"cmd.show('sphere', '{selection}')\n")
            file.write(f"cmd.color('white', '{residue_name}')\n")
            file.write(f"cmd.set('sphere_scale', 0.1, selection='{residue_name}')\n")
            file.write(f"cmd.zoom(selection = '{residue_name}', complete = 1)\n")

            for row in interaction_dataframe.itertuples(index=True):
                color=getattr(row, 'colors')
                # atom_name, residue_name, residue_number = row.Index.split("_")
                atom_name = row.Index.split("_")[0]
                # TODO this is hardcoded to 1:
                residue_number = '1'
                # residue_name, residue_number = atom_name[0][:3], atom_name[0][3:]
                if getattr(row, 'Probability') == .0: # This is the color white but in pymol is rendered in blue
                    # file.write(f"cmd.color('white', 'residue {residue_number} and name {atom_name[1]}')\n")
                    # file.write(f"cmd.set('sphere_scale', {pymol_vol_eq(getattr(row, 'Probability'), max_probability, min_probability)}, selection='residue {residue_number} and name {atom_name[1]}')\n")
                    pass
                else:
                    file.write(f"cmd.show('sphere', '{residue_name} and name {atom_name}')\n")
                    file.write(f"cmd.set('sphere_transparency', 0.5, selection='{residue_name}')\n")
                    file.write(f"cmd.set('sphere_scale', {pymol_vol_eq(getattr(row, 'Probability')/10, max_probability, min_probability)}, selection='{residue_name} and name {atom_name}')\n")
                    file.write(f"cmd.color('0x{color[1:]}', '{residue_name} and name {atom_name}')\n")
            file.write(f"cmd.set('ray_opaque_background', 1)\n")
            file.write(f"cmd.bg_color('white')\n")
            file.write(f"cmd.set('ray_shadows', 'off')\n")
            file.write(f"cmd.set('depth_cue', 0)\n")
            # file.write(f"cmd.set('ray_opaque_background', 'off')\n")
            # file.write(f"cmd.ray(1600, 800)\n")
            file.write(f"cmd.png('{image_output_folder}/{name}_{replica}_{trajectory}_{data}_pymol.png')\n")
        
        run_script(f'/home/emanuele/pymol-os/src/pymol_os/bin/pymol {simulation_output_folder}/{replica}/{name}_{replica}_{trajectory}_{data}.pml -c')

    return f'{image_output_folder}/{name}_{replica}_{trajectory}_{data}_pymol.png', f'{image_output_folder}/{name}_{replica}_{trajectory}_{data}_colorbar.png'


def pymol_vol_eq(probability, max_probability, min_probability):
    fraction = (0.5-0.2)/(max_probability-min_probability)
    scale = fraction*probability+0.3
    return scale


def merge_structure_colorbar(colorbar_image_path, structure_image_path, output_directory):
    img_structure = Image.open(structure_image_path).convert('RGBA')
    img_colorbar = Image.open(colorbar_image_path).convert('RGBA')
    img_colorbar = img_colorbar.rotate(90, expand=True)

    height1, width1 = img_structure.size
    height2, width2 = img_colorbar.size
    # Resize the images to have the same height
    # if height1 > height2:
    #     img_structure = img_structure.resize((height2, width1 * height2 // height1))
    # elif height2 > height1:
    # img_colorbar = img_colorbar.resize((height1, width2 * height1 // height2))

    # Crop the images to have the same width
    # if width1 > width2:
    #     img_structure = img_structure.crop((0, 0, width2, height2))
    # elif width2 > width1:
    #     img_colorbar = img_colorbar.crop((0, 0, width1, height1))

    img_colorbar = img_colorbar.resize((360, 480))

    full_img = Image.new('RGBA', (width1+360, height1))
    full_img.paste(img_colorbar, (width1,0), img_colorbar)
    full_img.paste(img_structure, (0,0))
    # full_img.save('./prova.png')

    datas = full_img.getdata()

    newData = []
    for item in datas:
        if item[0] == 255 and item[1] == 255 and item[2] == 255:
            newData.append((255, 255, 255, 0))
        else:
            newData.append(item)

    full_img.putdata(newData)
    full_img.save(output_directory, "PNG")


def make_contact_probability_comparison(contact_probability_fig, full_contact_probability_fig, simulation_name):
    compare_contact_probability_fig = go.Figure()
    compare_contact_probability_fig.add_trace(contact_probability_fig.data[2])
    compare_contact_probability_fig.update_traces(selector=dict(name='contact probability'), name=f'{simulation_name} covalent', line=dict(color=dm_tertiary_colors[1]))
    compare_contact_probability_fig.add_trace(full_contact_probability_fig[0].data[2])
    compare_contact_probability_fig.update_traces(selector=dict(name='contact probability'), name=f'{simulation_name} full', line=dict(color=dm_tertiary_colors[0]))
    
    for trace in compare_contact_probability_fig.data:
        trace.showlegend = True
    
    compare_contact_probability_fig.update_yaxes(
        title='Contact probability',
        # range=[0,1],
        showgrid=True,
        gridwidth=1,
        gridcolor='rgb(226, 226, 226)',
        showline=True,
        linecolor='black',
        linewidth=2,
        tickfont=dict(size=18),
        title_font=dict(size=22),
    )
    compare_contact_probability_fig.update_xaxes(
        title='',
        showline=True,
        linecolor='black',
        linewidth=2,
        tickfont=dict(size=18),
        title_font=dict(size=22),
    )
    compare_contact_probability_fig.update_layout(
        # title_text=title_plot,
        font=dict(size=title_font_size),
        plot_bgcolor='rgb(255,255,255)',
        legend=dict(
            orientation = 'h',
            yanchor='bottom',
            y=0.9,
            xanchor='right',
            x=1,
            font=dict(size=legend_font_size)
        )
    )

    return compare_contact_probability_fig


def get_DSSP_report(simulation_subset, ppt_presentation, report_folder):
    dssp_subset = filter_dict_by_string(simulation_subset, 'DSSP')
    temperature_subset = filter_dict_by_string(dssp_subset, 'full_temperature_trajectory')
    temperature_dssp_helix_fig = make_dssp_plot(temperature_subset, ['DSSP_helix_full'], report=True)
    temperature_dssp_sheet_fig = make_dssp_plot(temperature_subset, ['DSSP_sheet_full '], report=True)
    demux_subset = filter_dict_by_string(dssp_subset, 'demux')
    demux_dssp_helix_fig = make_dssp_plot(demux_subset, ['DSSP_helix'], report=True)
    demux_dssp_sheet_fig = make_dssp_plot(demux_subset, ['DSSP_sheet'], report=True)
    simulation_name = list(temperature_subset.keys())[0].split('|')[0]
    temperature_dssp_helix_fig.write_image(f'{report_folder}/png/{simulation_name}_replicas_temperature_dssp_helix_fig.png', format='png', scale=6, width=1000, height=500)
    temperature_dssp_sheet_fig.write_image(f'{report_folder}/png/{simulation_name}_replicas_temperature_dssp_sheet_fig.png', format='png', scale=6, width=1000, height=500)
    demux_dssp_helix_fig.write_image(f'{report_folder}/png/{simulation_name}_replicas_demux_dssp_helix_fig.png', format='png', scale=6, width=1000, height=500)
    demux_dssp_sheet_fig.write_image(f'{report_folder}/png/{simulation_name}_replicas_demux_dssp_sheet_fig.png', format='png', scale=6, width=1000, height=500)
    
    slide = ppt_presentation.slides.add_slide(ppt_presentation.slide_layouts[1])
    # slide.placeholders[0].text = f'{apo_simulation.plot_name}'
    slide.shapes.add_picture(f'{report_folder}/png/{simulation_name}_replicas_temperature_dssp_helix_fig.png',
                                left = Inches(1),
                                top = Inches(1),
                                width = Inches(12),
                                )
    
    slide = ppt_presentation.slides.add_slide(ppt_presentation.slide_layouts[1])
    slide.shapes.add_picture(f'{report_folder}/png/{simulation_name}_replicas_temperature_dssp_sheet_fig.png',
                                left = Inches(1),
                                top = Inches(1),
                                width = Inches(12),
                                )
    
    slide = ppt_presentation.slides.add_slide(ppt_presentation.slide_layouts[1])
    slide.shapes.add_picture(f'{report_folder}/png/{simulation_name}_replicas_demux_dssp_helix_fig.png',
                                left = Inches(1),
                                top = Inches(1),
                                width = Inches(12),
                                )
    
    slide = ppt_presentation.slides.add_slide(ppt_presentation.slide_layouts[1])
    slide.shapes.add_picture(f'{report_folder}/png/{simulation_name}_replicas_demux_dssp_sheet_fig.png',
                                left = Inches(1),
                                top = Inches(1),
                                width = Inches(12),
                                )
    # st.plotly_chart(demux_dssp_helix_fig)
    
    return ppt_presentation


def get_contact_matrix_report(simulation_subset, ppt_presentation, report_folder):
    contact_matrix_subset = filter_dict_by_string(simulation_subset, 'contact_matrix')
    temperature_subset = filter_dict_by_string(contact_matrix_subset, 'full_temperature_trajectory')

    for k in temperature_subset.keys():
        info = k.split('|')
        reduced_contact_matrix_subset = filter_dict_by_string(temperature_subset, k)
        contact_matrix_fig = make_contact_matrix_plot(reduced_contact_matrix_subset, report=True)
        contact_matrix_fig.write_image(f'{report_folder}/png/{info[0]}_{info[1]}_temperature_contact_matrix_fig.png', format='png', scale=4, width=500, height=500)
        del reduced_contact_matrix_subset

    num_plots = len(list(temperature_subset.keys()))
    num_rows = 2
    num_cols = 4
    num_slides = ceil(num_plots / (num_rows * num_cols))
    plot_index = 0

    for slide_num in range(num_slides):
        slide = ppt_presentation.slides.add_slide(ppt_presentation.slide_layouts[1])
        # plots_on_slide = min(num_rows * num_cols, num_plots - plot_index)
        
        # Arrange the plots in a grid on the current slide
        for row in range(num_rows):
            for col in range(num_cols):
                if plot_index < num_plots:
                    # plot = plots[plot_index]
                    plot = f'{report_folder}/png/{info[0]}_{plot_index}_temperature_contact_matrix_fig.png'
                    left = Inches(col * 3 + 0.8)
                    top = Inches(row * 3 + 1)
                    width = Inches(3)
                    height = Inches(3)
                    slide.shapes.add_picture(plot, left, top, width, height)
                    plot_index += 1
    
    demux_subset = filter_dict_by_string(contact_matrix_subset, 'demux')
    for k in demux_subset.keys():
        info = k.split('|')
        reduced_contact_matrix_subset = filter_dict_by_string(demux_subset, k)
        contact_matrix_fig = make_contact_matrix_plot(reduced_contact_matrix_subset, report=True)
        contact_matrix_fig.write_image(f'{report_folder}/png/{info[0]}_{info[1]}_demux_contact_matrix_fig.png', format='png', scale=4, width=500, height=500)
        del reduced_contact_matrix_subset

    num_plots = len(list(demux_subset.keys()))
    num_rows = 2
    num_cols = 4
    num_slides = ceil(num_plots / (num_rows * num_cols))
    plot_index = 0

    for slide_num in range(num_slides):
        slide = ppt_presentation.slides.add_slide(ppt_presentation.slide_layouts[1])
        # plots_on_slide = min(num_rows * num_cols, num_plots - plot_index)
        
        # Arrange the plots in a grid on the current slide
        for row in range(num_rows):
            for col in range(num_cols):
                if plot_index < num_plots:
                    # plot = plots[plot_index]
                    plot = f'{report_folder}/png/{info[0]}_{plot_index}_demux_contact_matrix_fig.png'
                    left = Inches(col * 3 + 0.8)
                    top = Inches(row * 3 + 1)
                    width = Inches(3)
                    height = Inches(3)
                    slide.shapes.add_picture(plot, left, top, width, height)
                    plot_index += 1
    
    return ppt_presentation


def get_gyr_salpha_report(simulation_subset, ppt_presentation, report_folder):
    fes_gyr_salpha_subset = filter_dict_by_string(simulation_subset, 'free_energy_gyration_salpha')
    temperature_subset = filter_dict_by_string(fes_gyr_salpha_subset, 'full_temperature_trajectory')

    for k in temperature_subset.keys():
        info = k.split('|')
        reduced_fes_gyr_salpha_subset = filter_dict_by_string(temperature_subset, k)
        contact_matrix_fig = make_fes_gyr_salpha_plot(reduced_fes_gyr_salpha_subset, report=True)
        contact_matrix_fig.write_image(f'{report_folder}/png/{info[0]}_{info[1]}_temperature_fes_gyr_salpha_fig.png', format='png', scale=4, width=500, height=500)
        del reduced_fes_gyr_salpha_subset

    num_plots = len(list(temperature_subset.keys()))
    num_rows = 2
    num_cols = 4
    num_slides = ceil(num_plots / (num_rows * num_cols))
    plot_index = 0

    for slide_num in range(num_slides):
        slide = ppt_presentation.slides.add_slide(ppt_presentation.slide_layouts[1])
        # plots_on_slide = min(num_rows * num_cols, num_plots - plot_index)
        
        # Arrange the plots in a grid on the current slide
        for row in range(num_rows):
            for col in range(num_cols):
                if plot_index < num_plots:
                    # plot = plots[plot_index]
                    plot = f'{report_folder}/png/{info[0]}_{plot_index}_temperature_fes_gyr_salpha_fig.png'
                    left = Inches(col * 3 + 0.8)
                    top = Inches(row * 3 + 1)
                    width = Inches(3)
                    height = Inches(3)
                    slide.shapes.add_picture(plot, left, top, width, height)
                    plot_index += 1

    demux_subset = filter_dict_by_string(fes_gyr_salpha_subset, 'demux')

    for k in demux_subset.keys():
        info = k.split('|')
        reduced_fes_gyr_salpha_subset = filter_dict_by_string(demux_subset, k)
        contact_matrix_fig = make_fes_gyr_salpha_plot(reduced_fes_gyr_salpha_subset, report=True)
        contact_matrix_fig.write_image(f'{report_folder}/png/{info[0]}_{info[1]}_demux_fes_gyr_salpha_fig.png', format='png', scale=4, width=500, height=500)
        del reduced_fes_gyr_salpha_subset

    num_plots = len(list(demux_subset.keys()))
    num_rows = 2
    num_cols = 4
    num_slides = ceil(num_plots / (num_rows * num_cols))
    plot_index = 0

    for slide_num in range(num_slides):
        slide = ppt_presentation.slides.add_slide(ppt_presentation.slide_layouts[1])
        # plots_on_slide = min(num_rows * num_cols, num_plots - plot_index)
        
        # Arrange the plots in a grid on the current slide
        for row in range(num_rows):
            for col in range(num_cols):
                if plot_index < num_plots:
                    # plot = plots[plot_index]
                    plot = f'{report_folder}/png/{info[0]}_{plot_index}_demux_fes_gyr_salpha_fig.png'
                    left = Inches(col * 3 + 0.8)
                    top = Inches(row * 3 + 1)
                    width = Inches(3)
                    height = Inches(3)
                    slide.shapes.add_picture(plot, left, top, width, height)
                    plot_index += 1
    
    return ppt_presentation


def get_header(summary_dict, name):
    header_str = f'''
{name}
{summary_dict[name]['Number of replicas']} replicas
Simulaton Length {summary_dict[name]['Simulation Length']}
Total Simulation Time {summary_dict[name]['Total Simulation Time']}
'''

    return header_str


def get_ligand_header(summary_dict, name, tsub):
    header_str = f'''
{name}
{summary_dict[name]['Number of replicas']} replicas
Simulaton Length {summary_dict[name]['Simulation Length']}
Total Simulation Time {summary_dict[name]['Total Simulation Time']}
Kd {summary_dict[name][f'{tsub.replace("_", " ")} Kd']}
'''

    return header_str


def replace_with_dict(input_string, replacement_dict):
    for old_str, new_str in replacement_dict.items():
        input_string = input_string.replace(old_str, new_str)
    return input_string


def get_structure_slides(simulation_subset, summary_dict, image_dict, ppt_presentation, report_folder, tsub, rename_dict):
    slide = ppt_presentation.slides.add_slide(ppt_presentation.slide_layouts[1])
    
    tsub_subset = filter_dict_by_string(simulation_subset, f'|0|{tsub}')
    simulation_name = list(tsub_subset.keys())[0].split('|')[0]
    name_header = replace_with_dict(simulation_name, rename_dict)
    if tsub == 'full_temperature_trajectory':
        tsub_header = ''
    else:
        tsub_header = tsub.replace("_", " ")
    
    slide.placeholders[0].text = f'{name_header} {tsub_header.replace("temperature ", "")}'
    textbox = slide.shapes.add_textbox(left = Inches(1),
                                        top = Inches(1),
                                        width = Inches(3.5),
                                        height = Inches(4)
                                        )
    text_frame = textbox.text_frame
    header_str = get_header(summary_dict, simulation_name)
    header_str = replace_with_dict(header_str, rename_dict)
    text_frame.text = header_str
    
    try:
        rdkit_img = image_dict[simulation_name]
        rdkit_image_stream = BytesIO()
        rdkit_img.save(rdkit_image_stream, format='PNG')
        rdkit_image_stream.seek(0)

        slide.shapes.add_picture(rdkit_image_stream,
            # left = Inches(6),
            # top = Inches(1.1),
            # width = Inches(4),
            left = Inches(6),
            top = Inches(1.1),
            width = Inches(3.8),
        )
    except: pass

    dssp_subset = filter_dict_by_string(tsub_subset, 'DSSP')
    # simulation_name = list(dssp_subset.keys())[0].split('|')[0]
    keys_iterator = iter(dssp_subset)
    first_key = next(keys_iterator)
    dssp_subset = {first_key:dssp_subset[first_key]}
    
    dssp_helix_fig = make_dssp_plot(dssp_subset, ['DSSP_helix', 'DSSP_sheet'])
    dssp_helix_fig.write_image(f'{report_folder}/png/{simulation_name}_{tsub}_dssp_helix_fig.png', format='png', scale=6, width=500, height=500)
    dssp_helix_fig.write_image(f'{report_folder}/svg/{simulation_name}_{tsub}_dssp_helix_fig.svg', format='svg', scale=6, width=500, height=500)
    dssp_helix_fig.write_html(f'{report_folder}/html/{simulation_name}_{tsub}_dssp_helix_fig.html')

    contact_matrix_subset = filter_dict_by_string(tsub_subset, 'contact_matrix')
    # simulation_name = list(contact_matrix_subset.keys())[0].split('|')[0]
    contact_matrix_fig = make_contact_matrix_plot(contact_matrix_subset)
    contact_matrix_fig.write_image(f'{report_folder}/png/{simulation_name}_{tsub}_contact_matrix_fig.png', format='png', scale=6, width=500, height=500)
    contact_matrix_fig.write_image(f'{report_folder}/svg/{simulation_name}_{tsub}_contact_matrix_fig.svg', format='svg', scale=6, width=500, height=500)
    contact_matrix_fig.write_html(f'{report_folder}/html/{simulation_name}_{tsub}_contact_matrix_fig.html')
    contact_matrix_fig = make_contact_matrix_plot(contact_matrix_subset)
    contact_matrix_fig.write_image(f'{report_folder}/png/{simulation_name}_{tsub}_contact_matrix_fig.png', format='png', scale=6, width=500, height=500)
    contact_matrix_fig.write_image(f'{report_folder}/svg/{simulation_name}_{tsub}_contact_matrix_fig.svg', format='svg', scale=6, width=500, height=500)
    contact_matrix_fig.write_html(f'{report_folder}/html/{simulation_name}_{tsub}_contact_matrix_fig.html')
    
    fes_gyr_salpha_subset = filter_dict_by_string(tsub_subset, 'free_energy_gyration_salpha')
    fes_gyr_salpha_fig = make_fes_gyr_salpha_plot(fes_gyr_salpha_subset)
    fes_gyr_salpha_fig.write_image(f'{report_folder}/png/{simulation_name}_{tsub}_fes_gyr_salpha_fig.png', format='png', scale=6, width=500, height=500)
    fes_gyr_salpha_fig.write_image(f'{report_folder}/svg/{simulation_name}_{tsub}_fes_gyr_salpha_fig.svg', format='svg', scale=6, width=500, height=500)
    fes_gyr_salpha_fig.write_html(f'{report_folder}/html/{simulation_name}_{tsub}_fes_gyr_salpha_fig.html')

    # TODO make 1D fes salpha
    fe_salpha_subset = filter_dict_by_string(tsub_subset, 'salpha_free_energy')
    fe_salpha_fig = make_salpha_plot(fe_salpha_subset)
    fe_salpha_fig.write_image(f'{report_folder}/png/{simulation_name}_{tsub}_fe_salpha_fig.png', format='png', scale=6, width=500, height=500)
    fe_salpha_fig.write_image(f'{report_folder}/svg/{simulation_name}_{tsub}_fe_salpha_fig.svg', format='svg', scale=6, width=500, height=500)
    fe_salpha_fig.write_html(f'{report_folder}/html/{simulation_name}_{tsub}_fe_salpha_fig.html')


    slide.shapes.add_picture(f'{report_folder}/png/{simulation_name}_{tsub}_dssp_helix_fig.png',
                                    left = Inches(1),
                                    top = Inches(3),
                                    width = Inches(3.82),
                                    )
    
    slide.shapes.add_picture(f'{report_folder}/png/{simulation_name}_{tsub}_contact_matrix_fig.png',
                                    left = Inches(5),
                                    top = Inches(3),
                                    width = Inches(3.5),
                                    )
    
    slide.shapes.add_picture(f'{report_folder}/png/{simulation_name}_{tsub}_fes_gyr_salpha_fig.png',
                                    left = Inches(9),
                                    top = Inches(3),
                                    width = Inches(3.33),
                                    )

    # slide.shapes.add_picture(f'{report_folder}/png/{simulation_name}_{tsub}_fe_salpha_fig.png',
    #                                 left = Inches(9),
    #                                 top = Inches(3),
    #                                 width = Inches(3.33),
    #                                 )

    return ppt_presentation, dssp_helix_fig


def get_ligand_slides(simulation_subset, summary_dict, image_dict, ppt_presentation, report_folder, tsub, apo_dssp_helix_fig, full_dssp_helix_fig, full_contact_probability_fig, rename_dict):
    # New slide with contact probability and Kd
    
    tsub_subset = filter_dict_by_string(simulation_subset, f'|0|{tsub}')
    simulation_name = list(tsub_subset.keys())[0].split('|')[0]
    simulation_output_folder = filter_dict_by_string(simulation_subset, 'full_temperature_trajectory|system_info')[f'{simulation_name}|0|full_temperature_trajectory|system_info']['output_folder'].to_string(index=False)
    residue_name = filter_dict_by_string(simulation_subset, 'full_temperature_trajectory|system_info')[f'{simulation_name}|0|full_temperature_trajectory|system_info']['ligand_name'].to_string(index=False)
    name_header = replace_with_dict(simulation_name, rename_dict)
    if tsub == 'full_temperature_trajectory':
        tsub_header = ''
    else:
        tsub_header = tsub.replace("_", " ")

    try:
        header_str = get_ligand_header(summary_dict, simulation_name, tsub)
    except:
        header_str = get_header(summary_dict, simulation_name)
        
    header_str = replace_with_dict(header_str, rename_dict)

    rdkit_img = image_dict[simulation_name]
    rdkit_image_stream = BytesIO()
    rdkit_img.save(rdkit_image_stream, format='PNG')
    rdkit_image_stream.seek(0)
    

    contact_probability = filter_dict_by_string(tsub_subset, 'ligand_contact_probability')
    keys_iterator = iter(contact_probability)
    first_key = next(keys_iterator)
    second_key = next(keys_iterator)
    contact_probability_fig = make_contact_probability_plot(contact_probability[first_key], ['contact_probability'])
    contact_probability_fig.write_image(f'{report_folder}/png/{simulation_name}_{tsub}_contact_probability_fig.png', format='png', scale=6, width=500, height=500)
    contact_probability_fig.write_image(f'{report_folder}/svg/{simulation_name}_{tsub}_contact_probability_fig.svg', format='svg', scale=6, width=500, height=500)
    contact_probability_fig.write_html(f'{report_folder}/html/{simulation_name}_{tsub}_contact_probability_fig.html')
    
    
    # Bound ensembles always have a Kd of 0
    if contact_probability[second_key]['Kd'].mean() != 0:
        slide = ppt_presentation.slides.add_slide(ppt_presentation.slide_layouts[1])
        slide.placeholders[0].text = f'{name_header} {tsub_header.replace("temperature ", "")}'
        textbox = slide.shapes.add_textbox(left = Inches(1),
                                            top = Inches(1),
                                            width = Inches(3.5),
                                            height = Inches(4)
                                            )
        text_frame = textbox.text_frame
        text_frame.text = header_str

        slide.shapes.add_picture(rdkit_image_stream,
            # left = Inches(6),
            # top = Inches(1.1),
            # width = Inches(4),
            left = Inches(6),
            top = Inches(1.1),
            width = Inches(3.8),
        )
        
        slide.shapes.add_picture(f'{report_folder}/png/{simulation_name}_{tsub}_contact_probability_fig.png',
                                        left = Inches(7),
                                        top = Inches(3),
                                        width = Inches(4.04),
                                        )
        
        kd_fig = make_kd_plot(contact_probability[second_key])
        kd_fig.write_image(f'{report_folder}/png/{simulation_name}_{tsub}_kd_fig.png', format='png', scale=6, width=800, height=500)
        kd_fig.write_image(f'{report_folder}/svg/{simulation_name}_{tsub}_kd_fig.svg', format='svg', scale=6, width=800, height=500)
        kd_fig.write_html(f'{report_folder}/html/{simulation_name}_{tsub}_kd_fig.html')
        slide.shapes.add_picture(f'{report_folder}/png/{simulation_name}_{tsub}_kd_fig.png',
                                        left = Inches(1),
                                        top = Inches(3),
                                        width = Inches(5.60),
                                        )


    hphob_contact_probability = filter_dict_by_string(tsub_subset, 'ligand_hydrophobic_contact_probability')
    keys_iterator = iter(hphob_contact_probability)
    first_key = next(keys_iterator)
    hphob_contact_probability_fig = make_contact_probability_plot(hphob_contact_probability[first_key], ['hydrophobic_contacts'], title_plot='Hydrophobic Contacts')
    hphob_contact_probability_fig.write_image(f'{report_folder}/png/{simulation_name}_{tsub}_hydrophobic_contact_probability_fig.png', format='png', scale=6, width=500, height=500)
    hphob_contact_probability_fig.write_image(f'{report_folder}/svg/{simulation_name}_{tsub}_hydrophobic_contact_probability_fig.svg', format='svg', scale=6, width=500, height=500)
    hphob_contact_probability_fig.write_html(f'{report_folder}/html/{simulation_name}_{tsub}_hydrophobic_contact_probability_fig.html')
    
    aromatic_contact_probability = filter_dict_by_string(tsub_subset, 'ligand_aromatic_contact_probability')    
    # The iterator is because there are two keys containing 'ligand_aromatic_contact_probability' and I need the first one only
    keys_iterator = iter(aromatic_contact_probability)
    first_key = next(keys_iterator)
    aromatic_contact_probability_fig = make_contact_probability_plot(aromatic_contact_probability[first_key], ['aromatic_stacking'], title_plot='Aromatic Stacking')
    aromatic_contact_probability_fig.write_image(f'{report_folder}/png/{simulation_name}_{tsub}_aromatic_contact_probability_fig.png', format='png', scale=6, width=500, height=500)
    aromatic_contact_probability_fig.write_image(f'{report_folder}/svg/{simulation_name}_{tsub}_aromatic_contact_probability_fig.svg', format='svg', scale=6, width=500, height=500)
    aromatic_contact_probability_fig.write_html(f'{report_folder}/html/{simulation_name}_{tsub}_aromatic_contact_probability_fig.html')
    
    try:
        hbond_contact_probability = filter_dict_by_string(tsub_subset, 'hbond_contact_probability')    
        # The iterator is because there are two keys containing 'ligand_aromatic_contact_probability' and I need the first one only
        keys_iterator = iter(hbond_contact_probability)
        first_key = next(keys_iterator)
        hbond_contact_probability_fig = make_contact_probability_plot(hbond_contact_probability[first_key], ['Hbonds_average', 'Hbonds_PD_average', 'Hbonds_LD_average'], title_plot='', show_legend=True)
        hbond_contact_probability_fig.write_image(f'{report_folder}/png/{simulation_name}_{tsub}_hbond_contact_probability_fig.png', format='png', scale=6, width=500, height=500)
        hbond_contact_probability_fig.write_image(f'{report_folder}/svg/{simulation_name}_{tsub}_hbond_contact_probability_fig.svg', format='svg', scale=6, width=500, height=500)
        hbond_contact_probability_fig.write_html(f'{report_folder}/html/{simulation_name}_{tsub}_hbond_contact_probability_fig.html')
    except: hbond_contact_probability_fig = None
    
    contact_probabilities_comparison_fig = make_ligand_contacts_comparison(hphob_contact_probability_fig, aromatic_contact_probability_fig, hbond_contact_probability_fig)
    contact_probabilities_comparison_fig.write_image(f'{report_folder}/png/{simulation_name}_{tsub}_contact_probabilities_comparison_fig.png', format='png', scale=6, width=500, height=500)
    contact_probabilities_comparison_fig.write_image(f'{report_folder}/svg/{simulation_name}_{tsub}_contact_probabilities_comparison_fig.svg', format='svg', scale=6, width=500, height=500)
    contact_probabilities_comparison_fig.write_html(f'{report_folder}/html/{simulation_name}_{tsub}_contact_probabilities_comparison_fig.html')
    
    dual_contact_probability = filter_dict_by_string(tsub_subset, 'ligand_dual_contact_probability')
    dual_contact_probability_fig = make_contact_matrix_plot(dual_contact_probability, dual='Dual ')
    dual_contact_probability_fig.write_image(f'{report_folder}/png/{simulation_name}_{tsub}_dual_contact_probability_fig.png', format='png', scale=6, width=500, height=500)
    dual_contact_probability_fig.write_image(f'{report_folder}/svg/{simulation_name}_{tsub}_dual_contact_probability_fig.svg', format='svg', scale=6, width=500, height=500)
    dual_contact_probability_fig.write_html(f'{report_folder}/html/{simulation_name}_{tsub}_dual_contact_probability_fig.html')


    # TODO on ligand stuff
    slide = ppt_presentation.slides.add_slide(ppt_presentation.slide_layouts[1])
    slide.placeholders[0].text = f'{name_header} {tsub_header.replace("temperature ", "")}'
    slide.shapes.add_picture(rdkit_image_stream,
        # left = Inches(6),
        # top = Inches(1.1),
        # width = Inches(4),
        left = Inches(1),
        top = Inches(1),
        width = Inches(3.8),
    )

    slide.shapes.add_picture(
        f'{report_folder}/png/{simulation_name}_{tsub}_contact_probabilities_comparison_fig.png',
        left = Inches(1),
        top = Inches(3),
        width = Inches(3.50),
    )
    
    any_onLigand_df = filter_dict_by_string(tsub_subset, 'any_contact_onLigand')
    # make_pymol_images(any_onLigand_df, residue_name, simulation_output_folder, max_probability=1, min_probability=0, image_output_folder = f'{report_folder}/png/{simulation_name}_{tsub}_all_contacts_pymol.png')
    any_contact_pymol_image_path, any_contact_pymol_colorbar_path = make_pymol_images(
        title='All Contacts',
        contact_dataframe=any_onLigand_df,
        residue_name=residue_name,
        simulation_output_folder=simulation_output_folder,
        image_output_folder=f'{report_folder}/png/')
    merge_structure_colorbar(any_contact_pymol_colorbar_path, any_contact_pymol_image_path, f'{report_folder}/png/{simulation_name}_{tsub}_all_contacts_pymol.png')

    any_hphob_onLigand_df = filter_dict_by_string(tsub_subset, 'any_hphob_contact_onLigand')
    any_hphob_contact_pymol_image_path, any_hphob_contact_pymol_colorbar_path = make_pymol_images(
        title='Hphobs',
        contact_dataframe=any_hphob_onLigand_df,
        residue_name=residue_name,
        simulation_output_folder=simulation_output_folder,
        image_output_folder=f'{report_folder}/png/')
    merge_structure_colorbar(any_hphob_contact_pymol_colorbar_path, any_hphob_contact_pymol_image_path, f'{report_folder}/png/{simulation_name}_{tsub}_hydrophobics_pymol.png')
    
    any_aromatics_onLigand_df = filter_dict_by_string(tsub_subset, 'onLigand_aromatic_contact_probability')    
    any_aromatics_contact_pymol_image_path, any_aromatics_contact_pymol_colorbar_path = make_pymol_images(
        title='Aro',
        contact_dataframe=any_aromatics_onLigand_df,
        residue_name=residue_name,
        simulation_output_folder=simulation_output_folder,
        image_output_folder=f'{report_folder}/png/')
    merge_structure_colorbar(any_aromatics_contact_pymol_colorbar_path, any_aromatics_contact_pymol_image_path, f'{report_folder}/png/{simulation_name}_{tsub}_aromatics_pymol.png')
    
    hbonds_onLigand_dict = filter_dict_by_string(tsub_subset, 'onLigand_hbonds_probability')
    for k, v in hbonds_onLigand_dict.items():
        try:
            df_temp = pd.DataFrame()
            df_temp['Probability'] = v['LD']
            hbonds_onLigand_split_dict = {f'{k}_LD' : df_temp}
            hbonds_LD_contact_pymol_image_path, hbonds_LD_contact_pymol_colorbar_path = make_pymol_images(
                title='HBonds_LD',
                contact_dataframe=hbonds_onLigand_split_dict,
                residue_name=residue_name,
                simulation_output_folder=simulation_output_folder,
                image_output_folder=f'{report_folder}/png/')
            merge_structure_colorbar(hbonds_LD_contact_pymol_colorbar_path, hbonds_LD_contact_pymol_image_path, f'{report_folder}/png/{simulation_name}_{tsub}_hbonds_LD_pymol.png')
        except: pass

        df_temp = pd.DataFrame()
        df_temp['Probability'] = v['PD']
        hbonds_onLigand_split_dict = {f'{k}_PD' : df_temp}
        hbonds_PD_contact_pymol_image_path, hbonds_PD_contact_pymol_colorbar_path = make_pymol_images(
            title='HBonds_PD',
            contact_dataframe=hbonds_onLigand_split_dict,
            residue_name=residue_name,
            simulation_output_folder=simulation_output_folder,
            image_output_folder=f'{report_folder}/png/')
        merge_structure_colorbar(hbonds_PD_contact_pymol_colorbar_path, hbonds_PD_contact_pymol_image_path, f'{report_folder}/png/{simulation_name}_{tsub}_hbonds_PD_pymol.png')

    slide.shapes.add_picture(
        f'{report_folder}/png/{simulation_name}_{tsub}_all_contacts_pymol.png',
        left = Inches(8.75),
        top = Inches(0.45),
        width = Inches(4.44),
    )
    slide.shapes.add_picture(
        f'{report_folder}/png/{simulation_name}_{tsub}_aromatics_pymol.png',
        left = Inches(4.50),
        top = Inches(2.50),
        width = Inches(4.44),
    )
    slide.shapes.add_picture(
        f'{report_folder}/png/{simulation_name}_{tsub}_hydrophobics_pymol.png',
        left = Inches(8.75),
        top = Inches(2.50),
        width = Inches(4.44),
    )
    slide.shapes.add_picture(
        f'{report_folder}/png/{simulation_name}_{tsub}_hbonds_PD_pymol.png',
        left = Inches(4.50),
        top = Inches(4.80),
        width = Inches(4.44),
    )
    try:
        slide.shapes.add_picture(
            f'{report_folder}/png/{simulation_name}_{tsub}_hbonds_LD_pymol.png',
            left = Inches(8.85),
            top = Inches(4.80),
            width = Inches(4.44),
        )
    except: pass

    # New slide with DSSP compare apo, dual contact and contact probability
    slide = ppt_presentation.slides.add_slide(ppt_presentation.slide_layouts[1])
    slide.placeholders[0].text = f'{name_header} {tsub_header.replace("temperature ", "")}'
    textbox = slide.shapes.add_textbox(left = Inches(1),
                                        top = Inches(1),
                                        width = Inches(3.5),
                                        height = Inches(4)
                                        )
    text_frame = textbox.text_frame
    text_frame.text = header_str
    
    slide.shapes.add_picture(rdkit_image_stream,
        # left = Inches(6),
        # top = Inches(1.1),
        # width = Inches(4),
        left = Inches(6),
        top = Inches(1.1),
        width = Inches(3.8),
    )

    # This one is to compare the DSSP with the apo simulation and the bound trajectory
    # It will return 3 lines
    # This is the current one
    dssp_subset = filter_dict_by_string(tsub_subset, 'DSSP')
    keys_iterator = iter(dssp_subset)
    first_key = next(keys_iterator)
    dssp_subset = {first_key:dssp_subset[first_key]}
    
    dssp_helix_fig = make_dssp_plot(dssp_subset, ['DSSP_helix'])
    dssp_helix_fig.update_traces(selector=dict(name='helix'), name=simulation_name, line=dict(color=dm_tertiary_colors[1]))
    # This the addition of apo
    if len(apo_dssp_helix_fig) != 0:
        dssp_helix_fig.add_trace(apo_dssp_helix_fig[0].data[0])
        dssp_helix_fig.update_traces(selector=dict(name='helix'), name='apo')
        # In case we would like to replace the labels.
        # dssp_helix_fig.update_traces(selector=dict(name='helix'), name=replace_with_dict('apo', rename_dict))
        dssp_helix_fig.update_layout(title_text='')
        
        # This is the addition of the bound trajectory
        if len(full_dssp_helix_fig) != 0:
            dssp_helix_fig.add_trace(full_dssp_helix_fig[0].data[0])
            # dssp_helix_fig.update_traces(selector=dict(name=simulation_name), name=f'{simulation_name} covalent', line=dict(color=dm_tertiary_colors[2]))
            dssp_helix_fig.update_traces(selector=dict(name=simulation_name), name=f'{simulation_name} {tsub.split("_")[0]}', line=dict(color=dm_tertiary_colors[2]))

        dssp_helix_fig.write_image(f'{report_folder}/png/{simulation_name}_{tsub}_dssp_compare_apo_fig.png', format='png', scale=6, width=500, height=500)
        dssp_helix_fig.write_image(f'{report_folder}/svg/{simulation_name}_{tsub}_dssp_compare_apo_fig.svg', format='svg', scale=6, width=500, height=500)
        dssp_helix_fig.write_html(f'{report_folder}/html/{simulation_name}_{tsub}_dssp_compare_apo_fig.html')
    
        slide.shapes.add_picture(f'{report_folder}/png/{simulation_name}_{tsub}_dssp_compare_apo_fig.png',
                                        left = Inches(1),
                                        top = Inches(3),
                                        width = Inches(3.82),
                                        )
    
    if 'full_temperature_trajectory' in tsub:
        dssp_helix_fig.update_traces(selector=dict(name=simulation_name), name=f'{simulation_name} full')

    # TODO keep in case Paul wants those separate.
    # if len(full_dssp_helix_fig) != 0:
    #     dssp_helix_fig.update_traces(selector=dict(name='helix'), name=simulation_name)
    #     dssp_helix_fig.add_trace(full_dssp_helix_fig[0].data[0])
    #     dssp_helix_fig.update_traces(selector=dict(name='helix'), name='full trajectory')
    
    #     dssp_helix_fig.write_image(f'{report_folder}/png/{simulation_name}_{tsub}_dssp_compare_full_fig.png', format='png', scale=6, width=500, height=500)
    #     dssp_helix_fig.write_image(f'{report_folder}/svg/{simulation_name}_{tsub}_dssp_compare_full_fig.svg', format='svg', scale=6, width=500, height=500)
    #     dssp_helix_fig.write_html(f'{report_folder}/html/{simulation_name}_{tsub}_dssp_compare_full_fig.html')
    
    #     slide.shapes.add_picture(f'{report_folder}/png/{simulation_name}_{tsub}_dssp_compare_apo_fig.png',
    #                                     left = Inches(1),
    #                                     top = Inches(3),
    #                                     width = Inches(3.50),
    #                                     )
    
    slide.shapes.add_picture(f'{report_folder}/png/{simulation_name}_{tsub}_dual_contact_probability_fig.png',
                                    left = Inches(5),
                                    top = Inches(3),
                                    width = Inches(3.50),
                                    )
    
    # This part is for adding a comparison between the contacts probabilities between the bound and the covalent bound
    # contact_probability_fig.update_traces(selector=dict(name='helix'), name=simulation_name, line=dict(color=dm_tertiary_colors[1]))
    if len(full_contact_probability_fig) != 0:
        compare_contact_probability_fig = make_contact_probability_comparison(contact_probability_fig, full_contact_probability_fig, simulation_name)
        compare_contact_probability_fig.write_image(f'{report_folder}/png/{simulation_name}_{tsub}_compare_contact_probability_fig.png', format='png', scale=6, width=500, height=500)
        compare_contact_probability_fig.write_image(f'{report_folder}/svg/{simulation_name}_{tsub}_compare_contact_probability_fig.svg', format='svg', scale=6, width=500, height=500)
        compare_contact_probability_fig.write_html(f'{report_folder}/html/{simulation_name}_{tsub}_compare_contact_probability_fig.html')
    

        slide.shapes.add_picture(f'{report_folder}/png/{simulation_name}_{tsub}_compare_contact_probability_fig.png',
                                        left = Inches(9),
                                        top = Inches(3),
                                        width = Inches(3.50),
                                        )
    else:
        slide.shapes.add_picture(f'{report_folder}/png/{simulation_name}_{tsub}_contact_probability_fig.png',
                                        left = Inches(9),
                                        top = Inches(3),
                                        width = Inches(3.50),
                                        )

    # New slide with contact probability, detailed contact probability and detailed hbonds
    slide = ppt_presentation.slides.add_slide(ppt_presentation.slide_layouts[1])
    slide.placeholders[0].text = f'{name_header} {tsub_header.replace("temperature ", "")}'
    textbox = slide.shapes.add_textbox(left = Inches(1),
                                        top = Inches(1),
                                        width = Inches(3.5),
                                        height = Inches(4)
                                        )
    
    text_frame = textbox.text_frame
    text_frame.text = header_str
    
    slide.shapes.add_picture(rdkit_image_stream,
        # left = Inches(6),
        # top = Inches(1.1),
        # width = Inches(4),
        left = Inches(6),
        top = Inches(1.1),
        width = Inches(3.8),
    )
    
    slide.shapes.add_picture(f'{report_folder}/png/{simulation_name}_{tsub}_contact_probability_fig.png',
                                    left = Inches(1),
                                    top = Inches(3),
                                    width = Inches(3.50),
                                    )

    slide.shapes.add_picture(f'{report_folder}/png/{simulation_name}_{tsub}_contact_probabilities_comparison_fig.png',
                                    left = Inches(5),
                                    top = Inches(3),
                                    width = Inches(3.50),
                                    )

    try:
        slide.shapes.add_picture(f'{report_folder}/png/{simulation_name}_{tsub}_hbond_contact_probability_fig.png',
                                        left = Inches(9),
                                        top = Inches(3),
                                        width = Inches(3.50),
                                        )
    except: pass

    return ppt_presentation, dssp_helix_fig, contact_probability_fig









