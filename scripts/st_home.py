import os
import datetime
import streamlit as st
from stqdm import stqdm
from pptx import Presentation
from itertools import product

from rest_simulation import REST_simulation
from st_report_functions import get_DSSP_report, get_contact_matrix_report, get_gyr_salpha_report, get_structure_slides, get_ligand_slides
from st_ligand import make_rdkit_image
from util import file_exists, filter_dict_by_string, st_get_simulation_info

from rdkit import Chem
from rdkit.Chem import AllChem
from rdkit.Chem import Draw
from rdkit.Chem.Draw import SimilarityMaps
from io import BytesIO
from PIL import Image


# @st.cache_data
def make_rdkit_report_image(smiles=None, output_folder=None):
    # if isinstance(smiles, str):
    #     mol = Chem.MolFromSmiles(smiles)
    # if isinstance(output_folder, str):
    #     mol = Chem.MolFromPDBFile(output_folder, removeHs=True)
    # mol = Chem.MolFromPDBFile(output_folder, removeHs=True)
    
    if smiles is not None:
        mol = Chem.MolFromSmiles(smiles)
    elif (output_folder is not None) and (smiles is None):
        mol = Chem.MolFromPDBFile(output_folder, removeHs=True)
    else:
        raise ValueError("Either 'smiles' or 'output_folder' must be provided.")

    # Ensure molecule is not None
    if mol is None:
        raise ValueError("Failed to create molecule from input.")

    AllChem.Compute2DCoords(mol)

    # # Save the molecule with 2D coordinates to a PDB file
    # if output_folder is not None:
    #     mol_pdb = Chem.AddHs(mol)
    #     AllChem.Compute2DCoords(mol_pdb)
    #     output_pdb_file = output_folder.replace(".pdb", "_2D.pdb")
    #     Chem.MolToPDBFile(mol_pdb, output_pdb_file)

    # Create a drawing object
    d = Draw.MolDraw2DCairo(800, 400)

    # Draw the molecule
    d.DrawMolecule(mol)

    # Finish drawing
    d.FinishDrawing()

    # Convert the drawing to an image
    img_data = BytesIO(d.GetDrawingText())
    img = Image.open(img_data).convert("RGB")

    # Display the image in Streamlit
    return img


# @st.cache_data
def save_dict(input_dict):
    temp_dict = {
        input_dict['name']:REST_simulation(name=input_dict['name'],
                                           rest_folder = input_dict['rest_folder'],
                                           sequence_offset = input_dict['sequence_offset'],
                                           output_folder = input_dict['output_folder'],
                                           number_of_replicas = int(input_dict['number_of_replicas']),
                                           replicas_toAnalyse = input_dict['replicas_toAnalyse']
                                           )
    }

    return temp_dict


# @st.cache_data
def get_simulation_list(new_simulations_dict):
    simulation_list = []
    for simulation in new_simulations_dict.keys():
        sim_name = simulation.split('|')[0]
        if sim_name not in simulation_list:
            simulation_list.append(sim_name)
    
    return simulation_list


def get_demux_report(new_simulations_dict, summary_dict, image_dict, selected_simulations, rename_dict, rest_summary, report_summary):
      
    ppt_rest_report = Presentation('../NUAGE_template.pptx')
    title_slide = ppt_rest_report.slide_layouts[0]
    slide = ppt_rest_report.slides.add_slide(title_slide)
    slide.placeholders[0].text = "REST replica report"
    slide.placeholders[1].text = """

    """

    ppt_simulation_report = Presentation('../NUAGE_template.pptx')
    title_slide = ppt_simulation_report.slide_layouts[0]
    slide = ppt_simulation_report.slides.add_slide(title_slide)
    slide.placeholders[0].text = "REST simulation report"
    slide.placeholders[1].text = """

    """

    current_working_directory = os.getcwd()

    report_folder = current_working_directory.replace('scripts', 'report/')
    if not os.path.exists(report_folder):
        os.makedirs(report_folder)
    
    report_folder_png = f'{report_folder}/png/'
    if not os.path.exists(report_folder_png):
        os.makedirs(report_folder_png)
    
    report_folder_svg = f'{report_folder}/svg/'
    if not os.path.exists(report_folder_svg):
        os.makedirs(report_folder_svg)
    
    report_folder_html = f'{report_folder}/html/'
    if not os.path.exists(report_folder_html):
        os.makedirs(report_folder_html)

    # TODO here add all the appends for the comparisons
    apo_dssp_helix = []
    for simulation in stqdm(selected_simulations):
        simulation_subset = filter_dict_by_string(new_simulations_dict, f'{simulation}|')
        trajectory_subset_keys = []
        for k in simulation_subset.keys():
            n = k.split('|')[2]
            if n not in trajectory_subset_keys and 'demux' not in n:
                trajectory_subset_keys.append(n)

        # TODO find a more clever way to define the position of the subsets
        # This is hardcoded assuming that covalent_bound comes before bound_trajectory
        if len(trajectory_subset_keys) == 3:
            trajectory_subset_keys = [trajectory_subset_keys[0], trajectory_subset_keys[2], trajectory_subset_keys[1]]

        if rest_summary is True:
            # DSSP
            ppt_rest_report = get_DSSP_report(simulation_subset, ppt_rest_report, report_folder)
            # Contact Matrix
            ppt_rest_report = get_contact_matrix_report(simulation_subset, ppt_rest_report, report_folder)
            # Gyration and Salpha
            ppt_rest_report = get_gyr_salpha_report(simulation_subset, ppt_rest_report, report_folder)

        if report_summary is True:
            full_dssp_helix, full_contact_probability = [], []
            for tsub in trajectory_subset_keys:
                if 'apo' in simulation:
                    ppt_simulation_report, dssp_comparison_fig = get_structure_slides(simulation_subset, summary_dict, image_dict, ppt_simulation_report, report_folder, tsub, rename_dict)
                    apo_dssp_helix.append(dssp_comparison_fig)
                else:
                # if 'apo' not in simulation:
                    ppt_simulation_report, _ = get_structure_slides(simulation_subset, summary_dict, image_dict, ppt_simulation_report, report_folder, tsub, rename_dict)
                    if 'full_temperature_trajectory' in tsub:
                        ppt_simulation_report, full_dssp_helix_fig, full_contact_probability_fig = get_ligand_slides(simulation_subset, summary_dict, image_dict, ppt_simulation_report, report_folder, tsub, apo_dssp_helix, full_dssp_helix, full_contact_probability, rename_dict)
                        full_dssp_helix.append(full_dssp_helix_fig)
                        full_contact_probability.append(full_contact_probability_fig)
                    else:
                        ppt_simulation_report, _, _ = get_ligand_slides(simulation_subset, summary_dict, image_dict, ppt_simulation_report, report_folder, tsub, apo_dssp_helix, full_dssp_helix, full_contact_probability, rename_dict)

    now = datetime.datetime.now()
    st.write('Writing pptx')
    # ppt_presentation.save(f'../test_pptx/NUAGE_ligand_updates_{now.month}_{now.day}_{now.hour}:{now.minute}.pptx')
    ppt_rest_report.save(f'{report_folder}/REST_report_{now.month}_{now.day}_{now.hour}_{now.minute}.pptx')
    ppt_simulation_report.save(f'{report_folder}/NUAGE_report_{now.month}_{now.day}_{now.hour}_{now.minute}.pptx')
        
    # return simulation_subset
    # return ppt_presentation


# # @st.cache_data
def create_summary(new_simulations_dict):
    simulation_names, simulation_replicas, simulation_subsets, _ = st_get_simulation_info(selected_simulations_dict=new_simulations_dict)
    summary_dict = {}
    image_dict = {}
    simulation_subsets = [item for item in simulation_subsets if 'tSNE' not in item]
    for name in simulation_names[::-1]:
        # name_subset = filter_dict_by_string(new_simulations_dict, name)
        kd_dict = {}
        system_info_dict = {}
        for subset in simulation_subsets:
            try:
                system_info_dict[f'Number of replicas'] = f"{new_simulations_dict[f'{name}|0|{subset}|system_info']['number of replicas'][0]}"
                system_info_dict[f'Simulation Length'] = f"{new_simulations_dict[f'{name}|0|{subset}|system_info']['single_replica_time'][0]} us"
                system_info_dict[f'Total Simulation Time'] = f"{new_simulations_dict[f'{name}|0|{subset}|system_info']['all_replica_time'][0]} us"
            except: pass
            try: kd_dict[f'{subset.replace("_", " ")} Kd'] = f"{round(new_simulations_dict[f'{name}|0|{subset}|ligand_contact_probability']['Kd'][0], 4)} mM"
            except: pass
        try:
            try:
                image_dict[name] = make_rdkit_report_image(smiles=f"{new_simulations_dict[f'{name}|0|{subset}|system_info']['smiles'][0]}", output_folder=f"{new_simulations_dict[f'{name}|0|{subset}|system_info']['ligand_pdb_path'][0]}")
            except:
                image_dict[name] = make_rdkit_report_image(output_folder=f"{new_simulations_dict[f'{name}|0|{subset}|system_info']['ligand_pdb_path'][0]}")
        except:
            pass
    
        summary_dict[name] = {**system_info_dict, **kd_dict}
        
    return summary_dict, image_dict


def home(new_simulations_dict):
    st.title('REST results')
    st.divider()
    
    # simulation_list = get_simulation_list(new_simulations_dict)
    simulation_names, simulation_replicas, simulation_subsets, _ = st_get_simulation_info(selected_simulations_dict=new_simulations_dict)
    simulation_list = simulation_names
    if 'apo' in simulation_list:
        simulation_list.remove('apo')
        simulation_list.insert(0, 'apo')
    # selected_data = [f'{sim}|{rep}|{sub}' for sim, rep, sub in product(simulation_names, simulation_replicas, simulation_subsets)]
    # selected_data = [f'{sim}|{rep}' for sim, rep in product(simulation_names, simulation_replicas)]

    summary_col, report_col = st.columns([0.6, 0.4])
    # summary_col, report_col = st.columns([0.2, 0.8])
    
    with summary_col:
        st.header('Summary')
        summary_dict, image_dict = create_summary(new_simulations_dict)
        with st.container(border=True, height=800):
            for data, value in summary_dict.items():
                data_col, image_col = st.columns([0.6, 0.4])
                with data_col:
                    st.subheader(data)
                    st.write(value)
                    # st.write(f':red[{data}]:', value)
                with image_col:
                    try: st.image(image_dict[data], use_column_width=True)
                    except: pass

    with report_col:
        st.header('Create report')
        with st.container(border=True, height=800):
            all_rep_col, rest_col, other_col = st.columns(3)
            with all_rep_col:
                show_all_replicas = st.checkbox('All replicas', value=True)
            with rest_col:
                rest_summary = st.checkbox('Make REST summary', value=False)
            with other_col:
                report_summary = st.checkbox('Make Report', value=True)
            if show_all_replicas:
                selected_simulation_toReport = simulation_list
            else:
                selected_simulation_toReport = st.multiselect("Select Simulation", simulation_list)
            
            with st.expander(label='Custom names'):
                rename_dict = {}
                for simulation in selected_simulation_toReport:
                    coso = st.text_input(label=simulation)
                    if coso:
                        rename_dict[simulation] = coso
                st.write('The following names will be replaced:')
                st.write(rename_dict)
            with st.form('Make report', border=True):
                st.write('The following simulations are selected to be included in the report.')
                st.write(selected_simulation_toReport)

                submitted = st.form_submit_button('Create Report')
                if submitted:
                    get_demux_report(new_simulations_dict, summary_dict, image_dict, selected_simulation_toReport, rename_dict, rest_summary, report_summary)
                    st.write(f'Report created in')
                    st.code('/stocazzo')

    


# def home_next(new_simulations_dict):
#     st.title('List of simulations')
#     # Get the constructor (__init__) method
#     init_method = getattr(REST_simulation, '__init__')

#     # Get the parameters of the constructor
#     params = inspect.signature(init_method).parameters

#     # Print the argument names    
#     st.write('Insert simulation details here:')
    
#     text_input_col, inputs_col = st.columns(2)

#     with text_input_col:
#         input_dict = {
#             'name' : st.text_input(label='Name'),
#             'rest_folder' : st.text_input(label='Rest Folder'),
#             'sequence_offset' : st.text_input(label='Sequence Offset'),
#             'output_folder' : st.text_input(label='Output Folder'),
#             'number_of_replicas' : st.text_input(label='Number of Replicas'),
#             'replicas_toAnalyse' : st.text_input(label='Replicas to Analyse')
#         }
    
#     if st.button(label='Save'):
#         st.write(new_simulations_dict)
#         temp_dict = save_dict(input_dict)
#         st.write(temp_dict)

#         st.write(new_simulations_dict)
        
#     with inputs_col:
#         st.write(input_dict['name'])
#         st.write(new_simulations_dict)
    
#     print(new_simulations_dict)
    
    # input_dict = {}
    # input_cols = st.columns(len(params))
    # for index, class_input in enumerate(params, 0):
    #     if class_input != 'self':
    #         with input_cols[index-1]:
    #             input_dict[class_input] = st.text_input(label=class_input.replace('_', ' '))
    # with input_cols[-1]:
    #     if st.button(label='Save'):
    #         simulations_dict[input_dict['name']] = save_dict(input_dict)
    
    
    
    # st.write(input_dict)
    
            # st.write("Arguments:", class_input)
















